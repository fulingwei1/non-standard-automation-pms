"""
基础幻灯片构建器 - 提供基本的幻灯片创建功能
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .config import PresentationConfig


class BaseSlideBuilder:
    """基础幻灯片构建器"""

    def __init__(self, presentation: Presentation):
        self.prs = presentation
        self.config = PresentationConfig()

    def add_title_slide(self, title: str, subtitle: str = ""):
        """
        添加标题幻灯片

        Args:
            title: 主标题
            subtitle: 副标题（可选）

        Returns:
            创建的幻灯片对象
        """
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)

        # 深色背景
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.config.DARK_BLUE
        background.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = self.config.TITLE_FONT_SIZE
        p.font.bold = True
        p.font.color.rgb = self.config.WHITE
        p.alignment = PP_ALIGN.CENTER

        # 副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2), Inches(9), Inches(1)
            )
            tf2 = sub_box.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.size = self.config.SUBTITLE_FONT_SIZE
            p2.font.color.rgb = self.config.TECH_BLUE
            p2.alignment = PP_ALIGN.CENTER

        return slide

    def add_section_slide(self, section_title: str, section_subtitle: str = ""):
        """
        添加章节分隔页

        Args:
            section_title: 章节标题
            section_subtitle: 章节副标题（可选）

        Returns:
            创建的幻灯片对象
        """
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 深蓝背景
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.config.DARK_BLUE
        background.line.fill.background()

        # 装饰线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.3), Inches(1), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.config.TECH_BLUE
        line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.5), Inches(9), Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = section_title
        p.font.size = self.config.SECTION_TITLE_FONT_SIZE
        p.font.bold = True
        p.font.color.rgb = self.config.WHITE

        if section_subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.5), Inches(9), Inches(0.8)
            )
            tf2 = sub_box.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = section_subtitle
            p2.font.size = Pt(20)
            p2.font.color.rgb = self.config.TECH_BLUE

        return slide

    def _add_white_background(self, slide):
        """添加白色背景"""
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.config.WHITE
        background.line.fill.background()
        return background

    def _add_top_bar(self, slide):
        """添加顶部装饰条"""
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.config.TOP_BAR_HEIGHT
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = self.config.TECH_BLUE
        top_bar.line.fill.background()
        return top_bar

    def _add_slide_title(self, slide, title: str):
        """添加幻灯片标题"""
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = self.config.HEADING_FONT_SIZE
        p.font.bold = True
        p.font.color.rgb = self.config.DARK_BLUE
        return title_box

    def _add_page_number(self, slide, page_num: int):
        """添加页码"""
        page_box = slide.shapes.add_textbox(
            Inches(9), Inches(7), Inches(0.8), Inches(0.3)
        )
        pf = page_box.text_frame.paragraphs[0]
        pf.text = str(page_num)
        pf.font.size = self.config.PAGE_NUMBER_FONT_SIZE
        pf.font.color.rgb = self.config.SILVER
        pf.alignment = PP_ALIGN.RIGHT
        return page_box

"""
内容幻灯片构建器 - 用于创建包含文字内容的幻灯片
"""

from pptx.util import Inches, Pt

from .base_builder import BaseSlideBuilder


class ContentSlideBuilder(BaseSlideBuilder):
    """内容幻灯片构建器"""

    def add_content_slide(self, title: str, content_list: list, page_num: int = None):
        """
        添加内容幻灯片

        Args:
            title: 幻灯片标题
            content_list: 内容列表，每个元素可以是字符串或字典
                          字典格式: {"text": str, "size": int, "bold": bool, "color": RgbColor, "level": int}
            page_num: 页码（可选）

        Returns:
            创建的幻灯片对象
        """
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 添加背景和装饰
        self._add_white_background(slide)
        self._add_top_bar(slide)
        self._add_slide_title(slide, title)

        # 内容区域
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.3), Inches(9), Inches(5.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content_list):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            if isinstance(item, dict):
                p.text = item.get("text", "")
                p.font.size = Pt(item.get("size", 18))
                p.font.bold = item.get("bold", False)
                if item.get("color"):
                    p.font.color.rgb = item["color"]
                else:
                    p.font.color.rgb = self.config.DARK_BLUE
                p.level = item.get("level", 0)
            else:
                p.text = str(item)
                p.font.size = self.config.CONTENT_FONT_SIZE
                p.font.color.rgb = self.config.DARK_BLUE

            p.space_after = Pt(8)

        # 页码
        if page_num:
            self._add_page_number(slide, page_num)

        return slide

    def add_two_column_slide(
        self, title: str, left_content: list, right_content: list, page_num: int = None
    ):
        """
        添加两栏内容幻灯片

        Args:
            title: 幻灯片标题
            left_content: 左栏内容列表
            right_content: 右栏内容列表
            page_num: 页码（可选）

        Returns:
            创建的幻灯片对象
        """
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 添加背景和装饰
        self._add_white_background(slide)
        self._add_top_bar(slide)
        self._add_slide_title(slide, title)

        # 左栏
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.3), Inches(4.3), Inches(5.5)
        )
        self._fill_textbox(left_box, left_content)

        # 右栏
        right_box = slide.shapes.add_textbox(
            Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.5)
        )
        self._fill_textbox(right_box, right_content)

        # 页码
        if page_num:
            self._add_page_number(slide, page_num)

        return slide

    def _fill_textbox(self, textbox, content_list: list):
        """填充文本框内容"""
        tf = textbox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(content_list):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = str(item)
            p.font.size = Pt(16)
            p.font.color.rgb = self.config.DARK_BLUE
            p.space_after = Pt(6)

"""
PPT生成器主类 - 统一管理幻灯片生成流程
"""

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .base_builder import BaseSlideBuilder
from .config import PresentationConfig
from .content_builder import ContentSlideBuilder
from .table_builder import TableSlideBuilder


class PresentationGenerator:
    """PPT生成器主类"""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = PresentationConfig.SLIDE_WIDTH
        self.prs.slide_height = PresentationConfig.SLIDE_HEIGHT

        self.base_builder = BaseSlideBuilder(self.prs)
        self.content_builder = ContentSlideBuilder(self.prs)
        self.table_builder = TableSlideBuilder(self.prs)
        self.config = PresentationConfig()

    def create_cover_slide(self):
        """创建封面 - 第1页"""
        slide = self.base_builder.add_title_slide(
            "智能驱动 · 精准交付", "非标自动化测试设备全生命周期 AI 项目管理系统"
        )
        # 添加slogan
        slogan_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
        p = slogan_box.text_frame.paragraphs[0]
        p.text = "让复杂项目变得可控、可预测、可信任"
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = self.config.SILVER
        p.alignment = PP_ALIGN.CENTER

    def create_toc_slide(self):
        """创建目录 - 第2页"""
        self.content_builder.add_content_slide(
            "内容导览",
            [
                {"text": "一、行业洞察与挑战", "size": 24, "bold": True},
                {"text": "二、系统架构总览", "size": 24, "bold": True},
                {"text": "三、核心功能详解", "size": 24, "bold": True},
                {"text": "四、客户门户与体验", "size": 24, "bold": True},
                {"text": "五、AI能力与技术亮点", "size": 24, "bold": True},
                {"text": "六、客户价值与案例", "size": 24, "bold": True},
                {"text": "七、合作展望", "size": 24, "bold": True},
            ],
            page_num=2,
        )

    def create_part1_industry_insights(self):
        """创建第一部分：行业洞察与挑战（第3-10页）"""
        # 第3页 - 公司简介
        self.content_builder.add_content_slide(
            "关于我们",
            [
                {
                    "text": "核心信息",
                    "size": 22,
                    "bold": True,
                    "color": self.config.TECH_BLUE,
                },
                "• 专注非标自动化测试设备研发制造 15+ 年",
                "• 服务客户覆盖汽车电子、消费电子、半导体、新能源等行业",
                "• 累计交付 1000+ 台定制化测试设备",
                "• 拥有 50+ 项自主知识产权专利",
                "",
                {
                    "text": "核心能力",
                    "size": 22,
                    "bold": True,
                    "color": self.config.TECH_BLUE,
                },
                "• ICT/FCT/EOL 测试设备",
                "• 烧录设备与老化系统",
                "• AOI 视觉检测设备",
                "• 自动化组装线体",
            ],
            page_num=3,
        )

        # 第4页 - 行业背景
        self.content_builder.add_content_slide(
            "非标自动化行业的黄金时代",
            [
                {
                    "text": "市场趋势",
                    "size": 22,
                    "bold": True,
                    "color": self.config.TECH_BLUE,
                },
                "• 全球智能制造市场规模预计2025年突破 5000亿美元",
                "• 中国非标自动化设备市场年复合增长率 15%+",
                "• 客户对交付速度要求:从6个月压缩至3个月",
                "• 设备定制化程度持续提升,标准化率不足 30%",
                "",
                {
                    "text": "驱动因素",
                    "size": 22,
                    "bold": True,
                    "color": self.config.TECH_BLUE,
                },
                "• 新能源汽车产业爆发式增长",
                "• 消费电子产品迭代加速",
                "• 半导体国产化进程提速",
                "• 人力成本持续上升",
            ],
            page_num=4,
        )

        # 第5页 - 客户画像
        self.table_builder.add_table_slide(
            "我们的客户",
            ["行业", "典型设备需求", "关键诉求"],
            [
                ["汽车电子", "BMS测试、ECU烧录、EOL检测", "高可靠性、可追溯"],
                ["消费电子", "PCBA测试、整机功能测试", "快速交付、柔性换线"],
                ["半导体", "晶圆测试、封装测试", "超高精度、洁净环境"],
                ["新能源", "电池模组测试、电机EOL", "大电流、高安全"],
            ],
            page_num=5,
        )

        # 第6-10页：各种痛点分析
        self._create_painpoint_slides()

    def _create_painpoint_slides(self):
        """创建痛点相关幻灯片"""
        # 第6页 - 痛点引入
        self.content_builder.add_content_slide(
            '非标项目的"复杂性困局"',
            [
                {
                    "text": "核心矛盾",
                    "size": 22,
                    "bold": True,
                    "color": self.config.ORANGE,
                },
                "高度定制化需求 × 紧迫交付周期 × 多专业协同 = 管理复杂度指数级上升",
                "",
                {
                    "text": "典型场景",
                    "size": 22,
                    "bold": True,
                    "color": self.config.TECH_BLUE,
                },
                "一个中型非标项目通常涉及:",
                "• 3-5 个专业团队（机械/电气/软件/测试/采购）",
                "• 200-500 个任务节点",
                "• 50-100 次客户变更请求",
                "• 1000+ 个物料SKU",
            ],
            page_num=6,
        )

        # ... 其他痛点幻灯片（7-10页）可以类似方式添加

    def create_part2_solution_overview(self):
        """创建第二部分：解决方案总览（第11-20页）"""
        # 第11页 - 转折页
        self.base_builder.add_section_slide("第二部分", "解决方案总览")

        # 第12页 - 转折
        self.content_builder.add_content_slide(
            '从"救火"走向"掌控"',
            [
                {
                    "text": "我们深耕非标自动化行业 15+ 年,深刻理解这些痛点。",
                    "size": 22,
                },
                {
                    "text": "基于 1000+ 项目经验,我们打造了这套 AI 驱动的项目管理系统。",
                    "size": 22,
                },
                "",
                {
                    "text": "设计理念",
                    "size": 24,
                    "bold": True,
                    "color": self.config.TECH_BLUE,
                },
                "• 不是通用工具的简单移植",
                '• 而是为非标自动化"量身定制"',
                "• 融入行业 Know-How 与 AI 智能",
            ],
            page_num=12,
        )

        # ... 其他幻灯片（13-20页）

    def create_part3_core_features(self):
        """创建第三部分：核心功能详解（第21-50页）"""
        # 第21页 - 章节分隔
        self.base_builder.add_section_slide("第三部分", "核心功能详解")

        # ... 功能详解幻灯片（22-50页）

    def create_part4_customer_portal(self):
        """创建第四部分：客户门户与体验（第51-60页）"""
        self.base_builder.add_section_slide("第四部分", "客户门户与体验")
        # ... 客户门户相关幻灯片

    def create_part5_ai_capabilities(self):
        """创建第五部分：AI能力与技术亮点（第61-70页）"""
        self.base_builder.add_section_slide("第五部分", "AI能力与技术亮点")
        # ... AI能力相关幻灯片

    def create_part6_customer_value(self):
        """创建第六部分：客户价值与案例（第71-80页）"""
        self.base_builder.add_section_slide("第六部分", "客户价值与案例")
        # ... 案例相关幻灯片

    def create_part7_cooperation(self):
        """创建第七部分：合作展望（第81-85页）"""
        self.base_builder.add_section_slide("第七部分", "合作展望")
        # ... 合作展望幻灯片

    def generate(self, output_path: str = "完整PPT.pptx"):
        """
        生成完整PPT

        Args:
            output_path: 输出文件路径

        Returns:
            输出文件路径
        """
        print("🚀 开始生成PPT...")

        # 1. 封面和目录
        print("  ✓ 创建封面和目录")
        self.create_cover_slide()
        self.create_toc_slide()

        # 2. 第一部分：行业洞察
        print("  ✓ 创建第一部分：行业洞察与挑战")
        self.create_part1_industry_insights()

        # 3. 第二部分：解决方案
        print("  ✓ 创建第二部分：解决方案总览")
        self.create_part2_solution_overview()

        # 4. 第三部分：核心功能
        print("  ✓ 创建第三部分：核心功能详解")
        self.create_part3_core_features()

        # 5-7. 其他部分
        print("  ✓ 创建其他部分...")
        self.create_part4_customer_portal()
        self.create_part5_ai_capabilities()
        self.create_part6_customer_value()
        self.create_part7_cooperation()

        # 保存
        print(f"  ✓ 保存PPT到: {output_path}")
        self.prs.save(output_path)

        print(f"✅ PPT生成完成: {output_path}")
        return output_path

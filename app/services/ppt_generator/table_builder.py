"""
表格幻灯片构建器 - 用于创建包含表格的幻灯片
"""

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from .base_builder import BaseSlideBuilder


class TableSlideBuilder(BaseSlideBuilder):
    """表格幻灯片构建器"""

    def add_table_slide(
        self, title: str, headers: list, rows: list, page_num: int = None
    ):
        """
        添加表格幻灯片

        Args:
            title: 幻灯片标题
            headers: 表头列表
            rows: 数据行列表，每行是一个列表
            page_num: 页码（可选）

        Returns:
            创建的幻灯片对象
        """
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 添加背景和装饰
        self._add_white_background(slide)
        self._add_top_bar(slide)
        self._add_slide_title(slide, title)

        # 创建表格
        num_rows = len(rows) + 1
        num_cols = len(headers)

        table_width = Inches(9)
        table_height = Inches(min(5, 0.5 * num_rows))

        table = slide.shapes.add_table(
            num_rows, num_cols, Inches(0.5), Inches(1.5), table_width, table_height
        ).table

        # 设置列宽
        col_width = table_width / num_cols
        for col in table.columns:
            col.width = int(col_width)

        # 表头
        self._format_table_header(table, headers)

        # 数据行
        self._format_table_rows(table, rows)

        # 页码
        if page_num:
            self._add_page_number(slide, page_num)

        return slide

    def _format_table_header(self, table, headers: list):
        """格式化表格头部"""
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.config.DARK_BLUE
            para = cell.text_frame.paragraphs[0]
            para.font.size = self.config.TABLE_HEADER_FONT_SIZE
            para.font.bold = True
            para.font.color.rgb = self.config.WHITE
            para.alignment = PP_ALIGN.CENTER

    def _format_table_rows(self, table, rows: list):
        """格式化表格数据行"""
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_data)
                # 斑马纹 - 偶数行浅蓝色背景
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.config.LIGHT_BLUE
                para = cell.text_frame.paragraphs[0]
                para.font.size = self.config.TABLE_CELL_FONT_SIZE
                para.font.color.rgb = self.config.DARK_BLUE
                para.alignment = PP_ALIGN.CENTER

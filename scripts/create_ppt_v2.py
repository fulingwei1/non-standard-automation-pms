#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
非标自动化项目管理系统 - 完整85页PPT生成器
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# 颜色定义
DARK_BLUE = RGBColor(10, 22, 40)
TECH_BLUE = RGBColor(0, 212, 255)
SILVER = RGBColor(192, 192, 192)
ORANGE = RGBColor(255, 107, 53)
GREEN = RGBColor(0, 196, 140)
WHITE = RGBColor(255, 255, 255)
LIGHT_BLUE = RGBColor(230, 247, 255)

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = TECH_BLUE
        p2.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_list, page_num=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = TECH_BLUE
    top_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, dict):
            p.text = item.get('text', '')
            p.font.size = Pt(item.get('size', 18))
            p.font.bold = item.get('bold', False)
            if item.get('color'):
                p.font.color.rgb = item['color']
            else:
                p.font.color.rgb = DARK_BLUE
            p.level = item.get('level', 0)
        else:
            p.text = str(item)
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_BLUE

        p.space_after = Pt(8)

    if page_num:
        page_box = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
        pf = page_box.text_frame.paragraphs[0]
        pf.text = str(page_num)
        pf.font.size = Pt(12)
        pf.font.color.rgb = SILVER
        pf.alignment = PP_ALIGN.RIGHT

    return slide

def add_section_slide(prs, section_title, section_subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.3), Inches(1), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = TECH_BLUE
    line.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if section_subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.8))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = section_subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = TECH_BLUE

    return slide

def add_table_slide(prs, title, headers, rows, page_num=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = TECH_BLUE
    top_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_width = Inches(9)
    table_height = Inches(min(5, 0.5 * num_rows))

    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.5), table_width, table_height).table

    col_width = table_width / num_cols
    for col in table.columns:
        col.width = int(col_width)

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(14)
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER

    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.color.rgb = DARK_BLUE
            para.alignment = PP_ALIGN.CENTER

    if page_num:
        page_box = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
        pf = page_box.text_frame.paragraphs[0]
        pf.text = str(page_num)
        pf.font.size = Pt(12)
        pf.font.color.rgb = SILVER
        pf.alignment = PP_ALIGN.RIGHT

    return slide

def add_two_column_slide(prs, title, left_content, right_content, page_num=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = TECH_BLUE
    top_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.3), Inches(5.5))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    for i, item in enumerate(left_content):
        if i == 0:
            p = tf_left.paragraphs[0]
        else:
            p = tf_left.add_paragraph()
        p.text = str(item)
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(6)

    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.5))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    for i, item in enumerate(right_content):
        if i == 0:
            p = tf_right.paragraphs[0]
        else:
            p = tf_right.add_paragraph()
        p.text = str(item)
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(6)

    if page_num:
        page_box = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
        pf = page_box.text_frame.paragraphs[0]
        pf.text = str(page_num)
        pf.font.size = Pt(12)
        pf.font.color.rgb = SILVER
        pf.alignment = PP_ALIGN.RIGHT

    return slide

def create_full_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========== 第一部分：开篇与背景（第1-10页）==========

    # 第1页 - 封面
    slide = add_title_slide(prs, "智能驱动 · 精准交付", "非标自动化测试设备全生命周期 AI 项目管理系统")
    slogan_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
    p = slogan_box.text_frame.paragraphs[0]
    p.text = "让复杂项目变得可控、可预测、可信任"
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = SILVER
    p.alignment = PP_ALIGN.CENTER

    # 第2页 - 目录
    add_content_slide(prs, "内容导览", [
        {"text": "一、行业洞察与挑战", "size": 24, "bold": True},
        {"text": "二、系统架构总览", "size": 24, "bold": True},
        {"text": "三、核心功能详解", "size": 24, "bold": True},
        {"text": "四、客户门户与体验", "size": 24, "bold": True},
        {"text": "五、AI能力与技术亮点", "size": 24, "bold": True},
        {"text": "六、客户价值与案例", "size": 24, "bold": True},
        {"text": "七、合作展望", "size": 24, "bold": True},
    ], page_num=2)

    # 第3页 - 公司简介
    add_content_slide(prs, "关于我们", [
        {"text": "核心信息", "size": 22, "bold": True, "color": TECH_BLUE},
        "• 专注非标自动化测试设备研发制造 15+ 年",
        "• 服务客户覆盖汽车电子、消费电子、半导体、新能源等行业",
        "• 累计交付 1000+ 台定制化测试设备",
        "• 拥有 50+ 项自主知识产权专利",
        "",
        {"text": "核心能力", "size": 22, "bold": True, "color": TECH_BLUE},
        "• ICT/FCT/EOL 测试设备",
        "• 烧录设备与老化系统",
        "• AOI 视觉检测设备",
        "• 自动化组装线体",
    ], page_num=3)

    # 第4页 - 行业背景
    add_content_slide(prs, "非标自动化行业的黄金时代", [
        {"text": "市场趋势", "size": 22, "bold": True, "color": TECH_BLUE},
        "• 全球智能制造市场规模预计2025年突破 5000亿美元",
        "• 中国非标自动化设备市场年复合增长率 15%+",
        "• 客户对交付速度要求：从6个月压缩至3个月",
        "• 设备定制化程度持续提升，标准化率不足 30%",
        "",
        {"text": "驱动因素", "size": 22, "bold": True, "color": TECH_BLUE},
        "• 新能源汽车产业爆发式增长",
        "• 消费电子产品迭代加速",
        "• 半导体国产化进程提速",
        "• 人力成本持续上升",
    ], page_num=4)

    # 第5页 - 客户画像
    add_table_slide(prs, "我们的客户",
        ["行业", "典型设备需求", "关键诉求"],
        [
            ["汽车电子", "BMS测试、ECU烧录、EOL检测", "高可靠性、可追溯"],
            ["消费电子", "PCBA测试、整机功能测试", "快速交付、柔性换线"],
            ["半导体", "晶圆测试、封装测试", "超高精度、洁净环境"],
            ["新能源", "电池模组测试、电机EOL", "大电流、高安全"],
        ], page_num=5)

    # 第6页 - 痛点引入
    add_content_slide(prs, "非标项目的复杂性困局", [
        {"text": "核心矛盾", "size": 22, "bold": True, "color": ORANGE},
        "高度定制化需求 x 紧迫交付周期 x 多专业协同 = 管理复杂度指数级上升",
        "",
        {"text": "典型场景", "size": 22, "bold": True, "color": TECH_BLUE},
        "一个中型非标项目通常涉及：",
        "• 3-5 个专业团队（机械/电气/软件/测试/采购）",
        "• 200-500 个任务节点",
        "• 50-100 次客户变更请求",
        "• 1000+ 个物料SKU",
    ], page_num=6)

    # 第7页 - 痛点详解1
    add_content_slide(prs, "痛点1 需求变更频繁，版本管理混乱", [
        {"text": "现状描述", "size": 20, "bold": True, "color": TECH_BLUE},
        "• 客户在项目周期内平均发起 15-30 次需求变更",
        "• 变更信息散落在微信、邮件、会议纪要中",
        "• 工程师经常按旧图纸加工导致返工",
        "• 版本号混乱：Rev.A/Rev.B/Final/Final_V2...",
        "",
        {"text": "业务影响", "size": 20, "bold": True, "color": ORANGE},
        "! 返工成本占项目总成本 10-25%",
        "! 因版本错误导致的废料损失",
        "! 客户投诉：说好的功能怎么没有?",
    ], page_num=7)

    # 第8页 - 痛点详解2
    add_content_slide(prs, "痛点2 多团队协同效率低下", [
        {"text": "现状描述", "size": 20, "bold": True, "color": TECH_BLUE},
        "• 机械工程师修改了安装位，电气工程师不知道",
        "• 软件开发完成，才发现硬件接口变了",
        "• 采购提前下单，设计变更后物料报废",
        "• 项目经理每天花 3 小时对齐信息",
        "",
        {"text": "典型对话", "size": 20, "bold": True, "color": SILVER},
        "机械组：我们改了治具尺寸，群里说过了",
        "电气组：没看到啊，线束都做好了",
        "",
        {"text": "业务影响", "size": 20, "bold": True, "color": ORANGE},
        "! 70% 的返工源于信息不同步",
    ], page_num=8)

    # 第9页 - 痛点详解3
    add_content_slide(prs, "痛点3 进度不透明，风险后知后觉", [
        {"text": "现状描述", "size": 20, "bold": True, "color": TECH_BLUE},
        "• 项目状态依赖周会汇报，信息滞后 3-5 天",
        "• 完成 80% 可能持续两周不动",
        "• 关键路径任务延误未被及时发现",
        "• 客户突然来访，项目经理手忙脚乱准备资料",
        "",
        {"text": "典型场景", "size": 20, "bold": True, "color": SILVER},
        "周一汇报：一切正常",
        "周五发现：视觉算法验证卡住两周了",
        "",
        {"text": "业务影响", "size": 20, "bold": True, "color": ORANGE},
        "! 项目延期率高达 40-60%",
        "! 被动救火而非主动预防",
    ], page_num=9)

    # 第10页 - 痛点详解4
    add_content_slide(prs, "痛点4 FAT/SAT 问题追溯困难", [
        {"text": "现状描述", "size": 20, "bold": True, "color": TECH_BLUE},
        "• 验收测试发现问题，无法快速定位责任环节",
        "• 历史问题记录分散，无法形成知识沉淀",
        "• 相同问题在不同项目重复出现",
        "• 客户要求提供完整追溯报告，临时拼凑",
        "",
        {"text": "业务影响", "size": 20, "bold": True, "color": ORANGE},
        "! 验收周期延长 2-4 周",
        "! 客户对品质管控能力质疑",
        "! 经验无法传承，老员工离职知识流失",
    ], page_num=10)

    # ========== 第二部分：解决方案总览（第11-20页）==========

    add_section_slide(prs, "第二部分", "解决方案总览")

    add_content_slide(prs, "从救火走向掌控", [
        {"text": "我们深耕非标自动化行业 15+ 年，深刻理解这些痛点。", "size": 22},
        {"text": "基于 1000+ 项目经验，我们打造了这套 AI 驱动的项目管理系统。", "size": 22},
        "",
        {"text": "设计理念", "size": 24, "bold": True, "color": TECH_BLUE},
        "• 不是通用工具的简单移植",
        "• 而是为非标自动化量身定制",
        "• 融入行业 Know-How 与 AI 智能",
    ], page_num=12)

    add_content_slide(prs, "产品定位", [
        {"text": "一句话定义", "size": 22, "bold": True, "color": TECH_BLUE},
        "专为非标自动化测试设备打造的 AI 增强型项目协同平台，",
        "覆盖从 RFQ 到售后的全生命周期管理",
        "",
        {"text": "核心价值主张", "size": 22, "bold": True, "color": TECH_BLUE},
        "精准：每个任务、每个变更、每个问题可追溯",
        "高效：AI 辅助决策，减少重复劳动",
        "协同：打破部门墙，信息实时同步",
        "透明：项目状态一目了然，风险提前预警",
    ], page_num=13)

    add_content_slide(prs, "系统架构", [
        {"text": "四层架构", "size": 22, "bold": True, "color": TECH_BLUE},
        "",
        {"text": "客户层 Customer Layer", "size": 18, "bold": True},
        "   客户Portal | 移动端 | 第三方系统集成 | 数据大屏",
        "",
        {"text": "业务层 Business Layer", "size": 18, "bold": True},
        "   销售管理 | 项目管理 | 生产管理 | 采购管理 | 售后服务",
        "",
        {"text": "智能层 AI Layer", "size": 18, "bold": True},
        "   WBS智能拆解 | 风险预警 | 资源优化 | 知识图谱",
        "",
        {"text": "数据层 Data Layer", "size": 18, "bold": True},
        "   项目数据 | 物料数据 | 工时数据 | 质量数据",
    ], page_num=14)

    add_content_slide(prs, "覆盖项目全生命周期", [
        {"text": "九大阶段", "size": 22, "bold": True, "color": TECH_BLUE},
        "",
        "S1 需求进入 -> S2 方案设计 -> S3 采购备料 -> S4 加工制造 -> S5 装配调试",
        "",
        "S9 质保结项 <- S8 现场安装 <- S7 包装发运 <- S6 出厂验收(FAT)",
        "",
        {"text": "每个阶段的管控要点", "size": 20, "bold": True, "color": TECH_BLUE},
        "• S1 需求进入：需求规格书、需求冻结确认",
        "• S2 方案设计：3D设计图、BOM、设计评审通过",
        "• S3 采购备料：物料齐套、长周期件跟踪",
        "• S4 加工制造：加工件完成、外协质量验收",
        "• S5 装配调试：整机调试OK、调试问题闭环",
    ], page_num=15)

    add_content_slide(prs, "覆盖项目全生命周期（续）", [
        {"text": "后续阶段管控要点", "size": 22, "bold": True, "color": TECH_BLUE},
        "",
        "• S6 出厂验收：FAT报告、客户签字确认",
        "• S7 包装发运：发货单据、物流跟踪",
        "• S8 现场安装：SAT报告、现场问题处理",
        "• S9 质保结项：结项报告、尾款回收",
        "",
        {"text": "健康度状态", "size": 22, "bold": True, "color": TECH_BLUE},
        "H1 正常（绿色）- 项目进展顺利",
        "H2 有风险（黄色）- 需要关注",
        "H3 阻塞（红色）- 需要立即处理",
        "H4 已完结（灰色）- 项目已结束",
    ], page_num=16)

    add_content_slide(prs, "七大核心能力", [
        {"text": "1. 需求管理", "size": 20, "bold": True, "color": TECH_BLUE},
        "   客户需求全程追溯",
        {"text": "2. WBS智能拆解", "size": 20, "bold": True, "color": TECH_BLUE},
        "   AI自动生成任务树",
        {"text": "3. 多专业协同", "size": 20, "bold": True, "color": TECH_BLUE},
        "   机/电/软无缝衔接",
        {"text": "4. 风险预警", "size": 20, "bold": True, "color": TECH_BLUE},
        "   智能识别潜在风险",
        {"text": "5. 进度追踪", "size": 20, "bold": True, "color": TECH_BLUE},
        "   实时掌握项目状态",
        {"text": "6. 交付物管理", "size": 20, "bold": True, "color": TECH_BLUE},
        "   图纸/文档版本控制",
        {"text": "7. 客户门户", "size": 20, "bold": True, "color": TECH_BLUE},
        "   透明化客户体验",
    ], page_num=17)

    add_table_slide(prs, "告别 Excel + 微信的时代",
        ["维度", "传统方式", "本系统"],
        [
            ["项目规划", "Excel手工排期", "AI智能拆解"],
            ["进度跟踪", "周会汇报", "实时看板"],
            ["团队协同", "微信群+邮件", "统一平台"],
            ["变更管理", "口头通知", "流程审批"],
            ["风险管控", "靠经验", "AI预测"],
            ["客户沟通", "被动回复", "主动推送"],
            ["知识沉淀", "存个人电脑", "统一知识库"],
        ], page_num=18)

    add_table_slide(prs, "为什么不用 Jira / Trello / 飞书？",
        ["对比维度", "通用工具", "本系统"],
        [
            ["行业适配", "需大量配置", "开箱即用"],
            ["业务流程", "简单任务管理", "全流程覆盖"],
            ["物料管理", "无", "BOM/采购集成"],
            ["设计协同", "无", "CAD/PLM联动"],
            ["AI能力", "基础提醒", "WBS/风险预测"],
            ["客户界面", "无", "专属Portal"],
        ], page_num=19)

    add_two_column_slide(prs, "技术架构",
        [
            "前端层",
            "- React + TypeScript",
            "- 移动端 H5",
            "- 数据可视化",
            "",
            "后端层",
            "- FastAPI (Python)",
            "- 微服务架构",
            "- RESTful API",
        ],
        [
            "数据层",
            "- MySQL (业务数据)",
            "- Redis (缓存)",
            "- Elasticsearch (搜索)",
            "",
            "AI层",
            "- 大语言模型",
            "- 知识图谱",
            "- 机器学习",
        ], page_num=20)

    # ========== 第三部分：核心功能详解（第21-50页）==========

    add_section_slide(prs, "第三部分", "核心功能详解")

    add_content_slide(prs, "功能模块全景图", [
        {"text": "经营驾驶舱", "size": 24, "bold": True, "color": TECH_BLUE},
        "   管理层实时数据看板，一目了然",
        "",
        {"text": "五大业务模块", "size": 22, "bold": True, "color": TECH_BLUE},
        "• 销售管理 - 从线索到签单",
        "• 项目管理 - 从立项到交付",
        "• 生产管理 - 从图纸到产品",
        "• 采购管理 - 从需求到入库",
        "• 售后服务 - 从报修到满意",
        "",
        {"text": "基础支撑", "size": 22, "bold": True, "color": TECH_BLUE},
        "• 基础数据管理 | 系统管理",
    ], page_num=22)

    add_content_slide(prs, "销售管理模块", [
        {"text": "从线索到签单，商机不流失", "size": 20, "color": SILVER},
        "",
        {"text": "核心功能", "size": 22, "bold": True, "color": TECH_BLUE},
        "线索管理 - 多渠道线索统一收集",
        "客户360 - 客户信息全景画像",
        "商机跟进 - 销售漏斗可视化",
        "报价管理 - 智能报价单生成",
        "合同管理 - 电子签约，版本控制",
        "",
        {"text": "业务价值", "size": 22, "bold": True, "color": GREEN},
        "• 线索转化率提升 30%",
        "• 报价效率提升 50%",
    ], page_num=23)

    add_content_slide(prs, "线索管理", [
        {"text": "1. 多渠道收集", "size": 20, "bold": True, "color": TECH_BLUE},
        "   官网表单 | 展会名片扫描 | 销售手动录入 | 第三方平台",
        "",
        {"text": "2. 线索分配", "size": 20, "bold": True, "color": TECH_BLUE},
        "   按区域/行业自动分配 | 负载均衡 | 支持手动调整",
        "",
        {"text": "3. 跟进记录", "size": 20, "bold": True, "color": TECH_BLUE},
        "   通话记录自动同步 | 拜访签到定位 | 跟进提醒",
        "",
        {"text": "4. 转化分析", "size": 20, "bold": True, "color": TECH_BLUE},
        "   线索来源分析 | 转化漏斗统计 | 销售行为分析",
    ], page_num=24)

    add_two_column_slide(prs, "客户360视图",
        [
            "客户基本信息",
            "• 企业信息、联系人",
            "• 组织架构",
            "• 行业标签、规模等级",
            "• 信用评级",
            "",
            "历史交易",
            "• 成交项目列表",
            "• 累计交易金额",
            "• 回款情况",
        ],
        [
            "沟通记录",
            "• 拜访记录",
            "• 邮件往来",
            "• 会议纪要",
            "",
            "项目进展",
            "• 在执行项目",
            "• 项目满意度",
            "• 问题反馈",
        ], page_num=25)

    add_content_slide(prs, "商机管理", [
        {"text": "销售漏斗阶段", "size": 22, "bold": True, "color": TECH_BLUE},
        "线索 -> 需求确认 -> 方案报价 -> 商务谈判 -> 合同签订 -> 赢单/丢单",
        "",
        {"text": "功能特点", "size": 22, "bold": True, "color": TECH_BLUE},
        "可视化销售漏斗",
        "预测销售额",
        "预计成单日期",
        "赢单概率评估",
        "竞争对手分析",
        "",
        {"text": "管理价值", "size": 22, "bold": True, "color": GREEN},
        "销售预测准确率提升 40%",
    ], page_num=26)

    add_content_slide(prs, "智能报价管理", [
        {"text": "1. 报价模板库", "size": 20, "bold": True, "color": TECH_BLUE},
        "   标准设备报价模板 | 历史报价参考 | 成本基准线",
        "",
        {"text": "2. 智能报价生成", "size": 20, "bold": True, "color": TECH_BLUE},
        "   根据配置自动计算 | 利润率预警 | 一键生成PDF",
        "",
        {"text": "3. 版本管理", "size": 20, "bold": True, "color": TECH_BLUE},
        "   多版本对比 | 变更追溯 | 客户确认记录",
        "",
        {"text": "4. 审批流程", "size": 20, "bold": True, "color": TECH_BLUE},
        "   超折扣审批 | 多级审批链 | 移动端审批",
    ], page_num=27)

    add_content_slide(prs, "合同全生命周期管理", [
        {"text": "合同流程", "size": 22, "bold": True, "color": TECH_BLUE},
        "起草 -> 内部评审 -> 客户确认 -> 电子签约 -> 归档备案 -> 履约监控",
        "",
        {"text": "功能特点", "size": 22, "bold": True, "color": TECH_BLUE},
        "合同模板库",
        "条款风险检查",
        "电子签章",
        "合同档案库",
        "履约节点提醒",
        "回款计划跟踪",
        "",
        {"text": "业务价值", "size": 22, "bold": True, "color": GREEN},
        "合同审批周期缩短 60%",
    ], page_num=28)

    add_content_slide(prs, "项目管理模块", [
        {"text": "从立项到交付，全程可控", "size": 20, "color": SILVER},
        "",
        {"text": "核心功能", "size": 22, "bold": True, "color": TECH_BLUE},
        "立项管理 - 规范立项流程",
        "WBS管理 - AI智能任务拆解",
        "进度管理 - 实时甘特图",
        "风险管理 - 智能预警",
        "文档管理 - 版本可追溯",
        "验收管理 - FAT/SAT流程化",
        "",
        {"text": "业务价值", "size": 22, "bold": True, "color": GREEN},
        "项目交付周期缩短 25% | 延期项目减少 40%",
    ], page_num=29)

    add_content_slide(prs, "规范化立项流程", [
        {"text": "立项流程", "size": 22, "bold": True, "color": TECH_BLUE},
        "销售移交 -> 技术评审 -> 资源评估 -> 立项审批 -> 项目启动",
        "",
        {"text": "立项检查清单", "size": 22, "bold": True, "color": TECH_BLUE},
        "客户需求规格书完整性",
        "技术可行性评估通过",
        "预算成本核算完成",
        "资源(人力/设备)可用性确认",
        "交期承诺合理性评估",
        "风险预判与应对方案",
    ], page_num=30)

    add_table_slide(prs, "AI 驱动的 WBS 智能拆解",
        ["传统方式", "AI智能拆解"],
        [
            ["项目经理凭经验手工拆解", "输入需求，AI自动生成任务树"],
            ["容易遗漏关键任务", "基于知识库，全面覆盖"],
            ["工时估算靠感觉", "参考历史数据，精准估算"],
            ["新人上手慢", "模板复用，快速启动"],
        ], page_num=31)

    add_content_slide(prs, "AI WBS 拆解引擎", [
        {"text": "工作原理", "size": 22, "bold": True, "color": TECH_BLUE},
        "",
        "Step 1：NLP解析需求文本",
        "Step 2：识别关键子系统/模块",
        "Step 3：匹配知识图谱模板",
        "Step 4：生成WBS任务树",
        "Step 5：预估工时与里程碑",
        "",
        {"text": "拆解准确率", "size": 22, "bold": True, "color": GREEN},
        "标准设备：95%+ | 定制设备：85%+ | 支持人工调整",
    ], page_num=32)

    add_table_slide(prs, "行业专属 WBS 模板库",
        ["设备类型", "典型WBS节点", "模板数量"],
        [
            ["ICT测试设备", "治具设计/针床加工/上位机开发", "50+"],
            ["FCT测试设备", "测试方案/接口设计/测试脚本", "40+"],
            ["EOL测试设备", "工艺分析/通讯开发/MES对接", "30+"],
            ["烧录设备", "烧录方案/治具设计/烧录软件", "20+"],
            ["AOI视觉设备", "光源选型/算法开发/标定调试", "30+"],
            ["组装线体", "机构设计/电控设计/PLC编程", "50+"],
        ], page_num=33)

    add_content_slide(prs, "实时甘特图", [
        {"text": "1. 可视化排期", "size": 20, "bold": True, "color": TECH_BLUE},
        "   任务条拖拽调整 | 依赖关系连线 | 关键路径高亮",
        "",
        {"text": "2. 进度追踪", "size": 20, "bold": True, "color": TECH_BLUE},
        "   计划 vs 实际对比 | 进度偏差预警 | 里程碑达成状态",
        "",
        {"text": "3. 资源视图", "size": 20, "bold": True, "color": TECH_BLUE},
        "   人员负载热力图 | 资源冲突提示 | 跨项目资源协调",
        "",
        {"text": "4. 多视图切换", "size": 20, "bold": True, "color": TECH_BLUE},
        "   甘特图视图 | 看板视图 | 列表视图",
    ], page_num=34)

    add_content_slide(prs, "AI 风险雷达", [
        {"text": "风险识别维度", "size": 22, "bold": True, "color": TECH_BLUE},
        "",
        "        进度风险",
        "           |",
        "   资源风险 -+- 质量风险",
        "           |",
        "        成本风险",
        "",
        {"text": "四维风险实时监控", "size": 20, "bold": True, "color": TECH_BLUE},
        "• 进度风险：任务延期、里程碑达成率",
        "• 成本风险：预算执行率、成本偏差",
        "• 质量风险：问题关闭率、返工率",
        "• 资源风险：人员负载、技能匹配",
    ], page_num=35)

    add_table_slide(prs, "AI风险预警场景",
        ["风险类型", "触发条件", "AI建议"],
        [
            ["进度延误", "关键任务超期3天", "调配资源/调整计划"],
            ["物料短缺", "长周期件未下单", "立即启动采购"],
            ["资源冲突", "人员负载>120%", "任务重新分配"],
            ["质量风险", "同类问题重复出现", "根因分析/流程改进"],
            ["成本超支", "实际成本>预算110%", "成本复盘/变更控制"],
        ], page_num=36)

    add_content_slide(prs, "风险预警实战案例", [
        {"text": "案例场景", "size": 22, "bold": True, "color": TECH_BLUE},
        "项目：某新能源电池模组EOL测试设备",
        "阶段：S3采购备料",
        "",
        {"text": "AI发现的风险", "size": 22, "bold": True, "color": ORANGE},
        "! 高压接触器（型号：XX-500A）交期12周，当前未下单",
        "   该物料位于关键路径，将导致S5装配调试延迟3周",
        "",
        {"text": "AI建议", "size": 22, "bold": True, "color": TECH_BLUE},
        "1. 立即启动采购流程（优先级：紧急）",
        "2. 联系备选供应商询价（已匹配3家）",
        "3. 评估是否可替代为国产型号",
        "",
        {"text": "处理结果：当天完成采购下单，项目按期交付", "size": 20, "color": GREEN},
    ], page_num=37)

    add_content_slide(prs, "跨专业协同看板", [
        {"text": "三专业联动", "size": 22, "bold": True, "color": TECH_BLUE},
        "",
        "     机械组       电气组       软件组",
        "    3D结构设计   电气原理图   测试脚本开发",
        "    治具设计     接线图       HMI界面",
        "",
        {"text": "协同机制", "size": 22, "bold": True, "color": TECH_BLUE},
        "任务关联：一处变更，相关方收到通知",
        "文档关联：机械图纸/电气图/代码同一任务下管理",
        "@通知：问题讨论直接@相关工程师",
        "评审会议：在线发起多专业评审",
    ], page_num=38)

    add_content_slide(prs, "工程变更全流程管理", [
        {"text": "ECN (工程变更通知) 流程", "size": 22, "bold": True, "color": TECH_BLUE},
        "变更申请 -> 影响评估 -> 审批决策 -> 执行变更 -> 关闭验证",
        "",
        {"text": "变更类型", "size": 22, "bold": True, "color": TECH_BLUE},
        "• 设计变更（机械/电气/软件）",
        "• 物料变更（替代/升级）",
        "• 工艺变更",
        "• 规格变更",
        "• 计划变更",
        "",
        {"text": "管理价值", "size": 22, "bold": True, "color": GREEN},
        "变更有据可查 | 影响范围可评估 | 避免私自变更",
    ], page_num=39)

    add_table_slide(prs, "交付物与文档管理",
        ["文档类型", "典型内容", "版本控制"],
        [
            ["设计图纸", "3D模型、2D图纸、爆炸图", "自动版本号"],
            ["电气图纸", "原理图、接线图、IO表", "变更追溯"],
            ["软件代码", "测试脚本、HMI程序", "Git集成"],
            ["技术文档", "操作手册、维护手册", "模板化管理"],
            ["验收文档", "FAT/SAT报告、调试记录", "签字确认"],
        ], page_num=40)

    add_content_slide(prs, "出厂验收（FAT）管理", [
        {"text": "FAT流程", "size": 22, "bold": True, "color": TECH_BLUE},
        "验收准备 -> 功能测试 -> 精度测试 -> 问题处理 -> 客户签字 -> 发货放行",
        "",
        {"text": "验收模板", "size": 22, "bold": True, "color": TECH_BLUE},
        "• 标准验收项目库（可配置）",
        "• 验收标准（PASS/FAIL判定依据）",
        "• 检验表单（支持拍照、填数据）",
        "• 电子签批",
        "",
        {"text": "问题管理", "size": 22, "bold": True, "color": TECH_BLUE},
        "问题分类（A/B/C类）-> 责任人指派 -> 整改跟踪 -> 关闭确认",
    ], page_num=41)

    add_content_slide(prs, "现场安装（SAT）管理", [
        {"text": "SAT流程", "size": 22, "bold": True, "color": TECH_BLUE},
        "发货跟踪 -> 现场到达 -> 开箱检验 -> 安装调试 -> 联调测试 -> 客户验收",
        "",
        {"text": "移动端支持", "size": 22, "bold": True, "color": TECH_BLUE},
        "签到定位",
        "现场拍照",
        "问题记录",
        "即时沟通",
        "验收签字",
        "",
        {"text": "现场问题处理", "size": 22, "bold": True, "color": TECH_BLUE},
        "问题快速上报 | 远程技术支持 | 备件申请 | 解决方案记录",
    ], page_num=42)

    add_content_slide(prs, "项目健康度看板", [
        {"text": "项目整体健康度评分：85分", "size": 24, "bold": True, "color": TECH_BLUE},
        "",
        "进度健康度：82分",
        "  - 里程碑达成率：90%",
        "  - 任务按时完成率：78%",
        "",
        "成本健康度：88分",
        "  - 预算执行率：92%",
        "  - 成本偏差：+3%",
        "",
        "质量健康度：90分",
        "  - 问题关闭率：95%",
        "  - 返工率：2%",
    ], page_num=43)

    add_content_slide(prs, "生产管理模块", [
        {"text": "从图纸到产品，精益制造", "size": 20, "color": SILVER},
        "",
        {"text": "核心功能", "size": 22, "bold": True, "color": TECH_BLUE},
        "生产计划 - 排产与调度",
        "工单管理 - 任务派发与报工",
        "物料管理 - 齐套检查与领料",
        "车间管理 - 工位与产能",
        "生产看板 - 实时进度展示",
        "",
        {"text": "业务价值", "size": 22, "bold": True, "color": GREEN},
        "生产效率提升 20% | 物料齐套率 95%+",
    ], page_num=44)

    add_content_slide(prs, "智能生产排程", [
        {"text": "排程考虑因素", "size": 22, "bold": True, "color": TECH_BLUE},
        "• 项目交期优先级",
        "• 设备产能限制",
        "• 人员技能匹配",
        "• 物料齐套状态",
        "• 外协件到货时间",
        "",
        {"text": "排程优化", "size": 22, "bold": True, "color": TECH_BLUE},
        "• AI自动排程建议",
        "• 瓶颈资源识别",
        "• 多方案对比",
        "• 动态调整",
    ], page_num=45)

    add_content_slide(prs, "工单全程追踪", [
        {"text": "工单流程", "size": 22, "bold": True, "color": TECH_BLUE},
        "工单创建 -> 任务派发 -> 开工报工 -> 过程检验 -> 完工入库",
        "",
        {"text": "移动端报工", "size": 22, "bold": True, "color": TECH_BLUE},
        "• 扫码开工",
        "• 进度上报",
        "• 异常反馈",
        "• 完工确认",
        "",
        {"text": "工时统计", "size": 22, "bold": True, "color": TECH_BLUE},
        "实际工时 vs 标准工时 | 效率分析 | 人员绩效",
    ], page_num=46)

    add_table_slide(prs, "物料齐套检查",
        ["项目", "应到物料", "已到库", "齐套率", "状态"],
        [
            ["PJ001", "350", "340", "97%", "正常"],
            ["PJ002", "280", "210", "75%", "有风险"],
            ["PJ003", "420", "180", "43%", "阻塞"],
        ], page_num=47)

    add_content_slide(prs, "采购管理模块", [
        {"text": "降本增效，供应可靠", "size": 20, "color": SILVER},
        "",
        {"text": "核心功能", "size": 22, "bold": True, "color": TECH_BLUE},
        "采购需求 - BOM驱动采购",
        "采购订单 - 全流程管理",
        "供应商管理 - 评估与准入",
        "到货管理 - 收货与检验",
        "采购分析 - 成本与绩效",
        "",
        {"text": "业务价值", "size": 22, "bold": True, "color": GREEN},
        "采购周期缩短 30% | 采购成本降低 5-10%",
    ], page_num=48)

    add_content_slide(prs, "从BOM到采购需求", [
        {"text": "流程", "size": 22, "bold": True, "color": TECH_BLUE},
        "设计BOM -> 审核发布 -> 需求汇总 -> 询价比价 -> 采购下单",
        "",
        {"text": "自动化特点", "size": 22, "bold": True, "color": TECH_BLUE},
        "• BOM变更自动触发采购需求更新",
        "• 物料需求合并（跨项目）",
        "• 安全库存预警",
        "• 长周期件提前预警",
    ], page_num=49)

    add_table_slide(prs, "供应商全生命周期管理",
        ["维度", "内容"],
        [
            ["准入管理", "资质审核、现场考察、试用评估"],
            ["绩效评估", "质量、交期、价格、服务"],
            ["等级分类", "战略/优选/合格/限用/淘汰"],
            ["风险监控", "经营异常、质量问题、交期风险"],
        ], page_num=50)

    # ========== 第四部分：客户门户与体验（第51-60页）==========

    add_section_slide(prs, "第四部分", "客户门户与体验")

    add_content_slide(prs, "客户专属 Portal", [
        {"text": "透明可信，自助便捷", "size": 20, "color": SILVER},
        "",
        {"text": "核心价值", "size": 22, "bold": True, "color": TECH_BLUE},
        "项目进度随时可查",
        "交付物在线查看",
        "问题沟通便捷",
        "验收签批在线完成",
        "",
        {"text": "客户体验升级", "size": 22, "bold": True, "color": GREEN},
        "从天天追问到自助查询",
        "从被动等待到主动推送",
    ], page_num=52)

    add_content_slide(prs, "项目进度实时可见", [
        {"text": "客户可见内容", "size": 22, "bold": True, "color": TECH_BLUE},
        "",
        {"text": "1. 项目总览", "size": 20, "bold": True},
        "   项目当前阶段 | 整体进度百分比 | 预计交付日期 | 健康度状态",
        "",
        {"text": "2. 里程碑进展", "size": 20, "bold": True},
        "   已完成里程碑 | 进行中里程碑 | 待开始里程碑",
        "",
        {"text": "3. 关键任务", "size": 20, "bold": True},
        "   本周完成任务 | 下周计划任务 | 阻塞问题说明",
    ], page_num=53)

    add_content_slide(prs, "交付物在线查看", [
        {"text": "可查看内容", "size": 22, "bold": True, "color": TECH_BLUE},
        "设计图纸（3D/2D在线预览）",
        "测试视频",
        "测试报告",
        "验收文档",
        "操作手册",
        "",
        {"text": "权限控制", "size": 22, "bold": True, "color": TECH_BLUE},
        "• 按项目阶段开放",
        "• 按文档类型授权",
        "• 水印保护",
        "• 下载记录",
    ], page_num=54)

    add_content_slide(prs, "FAT 远程验收", [
        {"text": "远程验收功能", "size": 22, "bold": True, "color": TECH_BLUE},
        "",
        {"text": "1. 视频直播", "size": 20, "bold": True},
        "   多机位现场直播 | 画面清晰流畅 | 支持回放",
        "",
        {"text": "2. 实时互动", "size": 20, "bold": True},
        "   语音对讲 | 文字聊天 | 标注问题点",
        "",
        {"text": "3. 在线签批", "size": 20, "bold": True},
        "   验收项确认 | 问题记录 | 电子签名",
        "",
        {"text": "业务价值", "size": 22, "bold": True, "color": GREEN},
        "客户差旅成本降低 70% | 验收周期缩短",
    ], page_num=55)

    add_content_slide(prs, "问题反馈与跟踪", [
        {"text": "反馈渠道", "size": 22, "bold": True, "color": TECH_BLUE},
        "在线提交问题单",
        "图片/视频上传",
        "一键呼叫服务热线",
        "在线客服",
        "",
        {"text": "问题跟踪", "size": 22, "bold": True, "color": TECH_BLUE},
        "• 工单状态实时更新",
        "• 处理进度推送",
        "• 历史问题查询",
        "• 满意度评价",
    ], page_num=56)

    add_content_slide(prs, "自助服务知识库", [
        {"text": "自助内容", "size": 22, "bold": True, "color": TECH_BLUE},
        "设备操作指南",
        "常见问题FAQ",
        "简单故障自排",
        "操作视频教程",
        "",
        {"text": "智能推荐", "size": 22, "bold": True, "color": TECH_BLUE},
        "• 根据设备型号推荐",
        "• 根据历史问题推荐",
        "• AI智能问答",
    ], page_num=57)

    add_content_slide(prs, "设备运行数据报表", [
        {"text": "可查看报表", "size": 22, "bold": True, "color": TECH_BLUE},
        "设备OEE报表",
        "测试数据趋势",
        "故障统计分析",
        "备件消耗报表",
        "",
        {"text": "数据可视化", "size": 22, "bold": True, "color": TECH_BLUE},
        "• 自定义时间范围",
        "• 多维度筛选",
        "• 图表/表格切换",
        "• 导出Excel/PDF",
    ], page_num=58)

    add_two_column_slide(prs, "移动端随时掌控",
        [
            "项目团队功能",
            "任务查看与更新",
            "待办提醒推送",
            "现场问题上报",
            "移动审批",
            "即时沟通",
        ],
        [
            "客户功能",
            "项目进度查看",
            "文档在线查看",
            "问题提交",
            "在线签批",
        ], page_num=59)

    add_table_slide(prs, "PC + 移动 + 大屏 多端协同",
        ["终端", "使用场景", "典型用户"],
        [
            ["PC Web", "日常办公、深度操作", "项目经理、工程师"],
            ["移动端", "外出、车间、现场", "服务工程师"],
            ["数据大屏", "管理驾驶舱、展厅展示", "高管、客户参观"],
            ["平板", "现场验收、会议演示", "验收工程师"],
        ], page_num=60)

    # ========== 第五部分：AI能力与技术亮点（第61-70页）==========

    add_section_slide(prs, "第五部分", "AI能力与技术亮点")

    add_content_slide(prs, "AI 赋能智能项目管理", [
        {"text": "AI能力矩阵", "size": 22, "bold": True, "color": TECH_BLUE},
        "",
        "WBS智能拆解 | 风险预警引擎 | 资源优化引擎 | 智能问答",
        "",
        "知识图谱引擎 | NLP处理引擎 | 数据分析引擎 | 机器学习",
    ], page_num=62)

    add_content_slide(prs, "非标行业知识图谱", [
        {"text": "知识图谱内容", "size": 22, "bold": True, "color": TECH_BLUE},
        "",
        {"text": "1. 设备知识库", "size": 20, "bold": True},
        "   500+ 典型设备类型 | 10000+ 标准模块 | 零部件关系图谱",
        "",
        {"text": "2. 工艺知识库", "size": 20, "bold": True},
        "   测试工艺流程 | 调试经验总结 | 故障处理方案",
        "",
        {"text": "3. 物料知识库", "size": 20, "bold": True},
        "   物料属性信息 | 替代物料关系 | 供应商信息",
        "",
        {"text": "4. 项目经验库", "size": 20, "bold": True},
        "   历史项目数据 | 工时基准 | 风险案例",
    ], page_num=63)

    add_table_slide(prs, "AI 风险预测引擎",
        ["风险类型", "输入特征", "预测输出"],
        [
            ["进度风险", "任务依赖、历史延误率", "延误概率、延误天数"],
            ["成本风险", "物料价格波动、变更频率", "超支概率、超支金额"],
            ["质量风险", "历史问题、新技术占比", "问题发生概率"],
            ["资源风险", "人员负载、技能匹配度", "资源冲突概率"],
        ], page_num=64)

    add_two_column_slide(prs, "AI 智能推荐",
        [
            "物料推荐",
            "• 替代物料推荐",
            "• 供应商推荐",
            "• 价格预测",
            "",
            "人员推荐",
            "• 任务分配建议",
            "• 技能匹配度",
            "• 负载均衡",
        ],
        [
            "解决方案推荐",
            "• 故障处理方案",
            "• 相似案例参考",
            "• 最佳实践",
            "",
            "进度调整建议",
            "• 赶工方案",
            "• 资源调配",
            "• 并行优化",
        ], page_num=65)

    add_content_slide(prs, "AI 智能助手", [
        {"text": "问答示例", "size": 22, "bold": True, "color": TECH_BLUE},
        "",
        {"text": "用户问：PJ250101项目当前状态如何？", "size": 18, "color": SILVER},
        "",
        {"text": "AI回答：", "size": 18, "bold": True},
        "PJ250101项目当前处于S5装配调试阶段，",
        "整体进度75%，健康度评分82分（有风险）。",
        "主要风险：视觉算法调试进度落后3天。",
        "建议：增加调试人员或申请延期。",
    ], page_num=66)

    add_two_column_slide(prs, "开放集成生态",
        [
            "CAD/PLM集成",
            "• SolidWorks",
            "• AutoCAD",
            "• CATIA",
            "• Creo",
            "",
            "ERP集成",
            "• SAP",
            "• 金蝶",
            "• 用友",
        ],
        [
            "自动化软件",
            "• TIA Portal",
            "• RSLogix",
            "• NI TestStand",
            "• LabVIEW",
            "",
            "办公协作",
            "• 企业微信",
            "• 钉钉",
            "• 飞书",
        ], page_num=67)

    add_content_slide(prs, "数据分析与经营洞察", [
        {"text": "分析维度", "size": 22, "bold": True, "color": TECH_BLUE},
        "",
        {"text": "1. 项目分析", "size": 20, "bold": True},
        "   项目成功率趋势 | 延期原因分析 | 成本偏差分析",
        "",
        {"text": "2. 人效分析", "size": 20, "bold": True},
        "   人均产出 | 工时利用率 | 技能分布",
        "",
        {"text": "3. 供应链分析", "size": 20, "bold": True},
        "   供应商绩效 | 采购成本趋势 | 物料周转率",
        "",
        {"text": "4. 客户分析", "size": 20, "bold": True},
        "   客户满意度 | 复购率 | 客户价值分布",
    ], page_num=68)

    add_content_slide(prs, "高管经营驾驶舱", [
        {"text": "驾驶舱指标", "size": 22, "bold": True, "color": TECH_BLUE},
        "",
        "在执行项目: 25个 | 本月交付: 8个 | 在手订单: 5.2亿",
        "",
        "项目健康度分布",
        "正常: 18个   有风险: 5个   阻塞: 2个",
        "",
        "本月关键指标",
        "交付准时率: 88% | 客户满意度: 92% | 回款率: 78%",
        "",
        "风险预警",
        "PJ005 视觉算法调试延误",
        "PJ012 长周期件交期风险",
    ], page_num=69)

    add_table_slide(prs, "报表中心",
        ["报表类型", "报表名称", "周期"],
        [
            ["项目报表", "项目进度周报", "周"],
            ["项目报表", "里程碑达成报告", "月"],
            ["成本报表", "项目成本分析", "月"],
            ["人效报表", "工时利用率报告", "周"],
            ["采购报表", "采购执行月报", "月"],
            ["质量报表", "问题分析报告", "月"],
            ["客户报表", "客户满意度报告", "季"],
        ], page_num=70)

    # ========== 第六部分：客户价值与案例（第71-80页）==========

    add_section_slide(prs, "第六部分", "客户价值与案例")

    add_content_slide(prs, "为您创造的价值", [
        {"text": "交付更快", "size": 28, "bold": True, "color": TECH_BLUE},
        "项目交付周期缩短 25%",
        "",
        {"text": "成本更低", "size": 28, "bold": True, "color": TECH_BLUE},
        "项目变更成本降低 35%",
        "",
        {"text": "质量更优", "size": 28, "bold": True, "color": TECH_BLUE},
        "返工率降低至 5% 以下",
        "",
        {"text": "客户更满意", "size": 28, "bold": True, "color": TECH_BLUE},
        "客户满意度 NPS 提升至 82+",
    ], page_num=72)

    add_table_slide(prs, "交付效率提升",
        ["环节", "效率提升", "原因"],
        [
            ["项目启动", "50%", "AI WBS拆解，模板复用"],
            ["进度管控", "40%", "实时追踪，提前预警"],
            ["变更处理", "60%", "流程化管理，影响评估"],
            ["团队协同", "35%", "信息同步，减少会议"],
            ["验收交付", "30%", "在线验收，问题追踪"],
        ], page_num=73)

    add_content_slide(prs, "成本节约分析", [
        {"text": "成本节约来源", "size": 22, "bold": True, "color": TECH_BLUE},
        "",
        {"text": "1. 返工成本减少", "size": 20, "bold": True},
        "   信息同步减少版本错误 | 变更管理减少遗漏",
        "",
        {"text": "2. 人力成本优化", "size": 20, "bold": True},
        "   自动化报表减少人工统计 | 协同效率提升减少加班",
        "",
        {"text": "3. 物料成本控制", "size": 20, "bold": True},
        "   需求变更及时通知采购 | 减少废料损失",
        "",
        {"text": "ROI计算", "size": 22, "bold": True, "color": GREEN},
        "系统投入 vs 年化收益 = 1:5+",
    ], page_num=74)

    add_table_slide(prs, "风险可控，延期减少",
        ["指标", "使用前", "使用后", "改善"],
        [
            ["项目延期率", "45%", "15%", "降低30%"],
            ["重大风险发生", "月均3次", "月均0.5次", "降低83%"],
            ["风险发现时间", "事后", "提前2周", "-"],
            ["问题闭环周期", "5天", "2天", "降低60%"],
        ], page_num=75)

    add_table_slide(prs, "客户满意度提升",
        ["维度", "改善点", "客户反馈"],
        [
            ["透明度", "项目进度随时可查", "终于不用天天追问了"],
            ["响应速度", "问题2小时内响应", "比以前快多了"],
            ["沟通效率", "在线沟通留痕", "不怕说不清楚了"],
            ["交付质量", "问题追溯完整", "质量管控很到位"],
            ["服务体验", "自助查询方便", "用起来很顺手"],
        ], page_num=76)

    add_content_slide(prs, "案例一：某汽车电子企业", [
        {"text": "客户背景", "size": 22, "bold": True, "color": TECH_BLUE},
        "行业：新能源汽车电子 | 规模：年产能5000台设备",
        "挑战：项目多、变更频繁、交付压力大",
        "",
        {"text": "解决方案", "size": 22, "bold": True, "color": TECH_BLUE},
        "部署项目管理系统 | 对接现有ERP/PLM | 培训80+用户",
        "",
        {"text": "实施效果", "size": 22, "bold": True, "color": GREEN},
        "• 项目交付周期：从16周 -> 12周",
        "• 项目延期率：从40% -> 10%",
        "• 变更成本：降低35%",
        "• 客户满意度：从75 -> 90",
        "",
        {"text": "这套系统真正理解我们非标行业的痛点", "size": 18, "color": SILVER},
    ], page_num=77)

    add_content_slide(prs, "案例二：某消费电子测试服务商", [
        {"text": "客户背景", "size": 22, "bold": True, "color": TECH_BLUE},
        "行业：3C消费电子测试 | 规模：服务50+品牌客户",
        "挑战：项目数量大、客户要求高、进度跟踪难",
        "",
        {"text": "解决方案", "size": 22, "bold": True, "color": TECH_BLUE},
        "客户Portal上线 | 移动端报工 | AI风险预警",
        "",
        {"text": "实施效果", "size": 22, "bold": True, "color": GREEN},
        "• 同时在管项目：从30个 -> 50个",
        "• 项目经理效率：提升40%",
        "• 客户投诉：降低60%",
        "• 回款周期：缩短20天",
        "",
        {"text": "客户Portal让我们的专业度提升了一个台阶", "size": 18, "color": SILVER},
    ], page_num=78)

    add_content_slide(prs, "案例三：某半导体设备企业", [
        {"text": "客户背景", "size": 22, "bold": True, "color": TECH_BLUE},
        "行业：半导体封测设备 | 规模：高端定制设备供应商",
        "挑战：技术复杂、多方协同、知识沉淀难",
        "",
        {"text": "解决方案", "size": 22, "bold": True, "color": TECH_BLUE},
        "知识库建设 | 多专业协同看板 | 文档版本管理",
        "",
        {"text": "实施效果", "size": 22, "bold": True, "color": GREEN},
        "• 设计评审效率：提升50%",
        "• 图纸版本错误：降低90%",
        "• 新人上手周期：从3个月 -> 1个月",
        "• 知识复用率：从10% -> 60%",
        "",
        {"text": "终于不怕老员工离职了，经验都沉淀在系统里", "size": 18, "color": SILVER},
    ], page_num=79)

    add_content_slide(prs, "实施路径与服务", [
        {"text": "标准实施流程", "size": 22, "bold": True, "color": TECH_BLUE},
        "",
        "第1-2周        第3-4周        第5-6周        第7-8周",
        "需求调研 ->    系统配置 ->    数据迁移 ->    培训上线",
        "业务梳理      定制开发      试运行        运维支持",
        "",
        {"text": "实施内容", "size": 22, "bold": True, "color": TECH_BLUE},
        "• 需求调研：业务流程梳理、需求确认",
        "• 系统配置：参数配置、模板定制",
        "• 数据迁移：历史数据导入、基础数据",
        "• 培训上线：用户培训、试运行",
        "• 运维支持：问题处理、持续优化",
    ], page_num=80)

    # ========== 第七部分：合作展望（第81-85页）==========

    add_section_slide(prs, "第七部分", "合作展望")

    add_table_slide(prs, "产品版本与定价",
        ["功能", "标准版", "专业版", "企业版"],
        [
            ["项目管理", "有", "有", "有"],
            ["销售管理", "-", "有", "有"],
            ["生产管理", "-", "有", "有"],
            ["客户Portal", "-", "有", "有"],
            ["AI能力", "基础", "完整", "完整+定制"],
            ["系统集成", "-", "标准API", "深度集成"],
            ["部署方式", "SaaS", "SaaS/私有化", "私有化"],
            ["用户数", "20", "100", "不限"],
        ], page_num=82)

    add_content_slide(prs, "服务保障体系", [
        {"text": "1. 实施服务", "size": 20, "bold": True, "color": TECH_BLUE},
        "   专业实施顾问 | 项目经理负责制 | 定期项目汇报",
        "",
        {"text": "2. 培训服务", "size": 20, "bold": True, "color": TECH_BLUE},
        "   管理员培训 | 用户培训 | 在线培训课程",
        "",
        {"text": "3. 技术支持", "size": 20, "bold": True, "color": TECH_BLUE},
        "   7x24小时热线 | 远程技术支持 | 2小时响应承诺",
        "",
        {"text": "4. 持续服务", "size": 20, "bold": True, "color": TECH_BLUE},
        "   季度业务回顾 | 系统优化建议 | 新功能培训",
        "",
        {"text": "SLA承诺", "size": 22, "bold": True, "color": GREEN},
        "系统可用性：99.9% | 问题响应：2小时 | 重大问题解决：24小时",
    ], page_num=83)

    add_content_slide(prs, "开启您的数字化项目管理之旅", [
        {"text": "下一步", "size": 22, "bold": True, "color": TECH_BLUE},
        "",
        {"text": "1. 预约演示", "size": 22, "bold": True},
        "   30分钟系统演示 | 针对您的业务场景 | 解答您的疑问",
        "",
        {"text": "2. 试用体验", "size": 22, "bold": True},
        "   14天免费试用 | 真实数据测试 | 专人辅导",
        "",
        {"text": "3. 方案咨询", "size": 22, "bold": True},
        "   业务需求梳理 | 定制化方案 | 投资回报分析",
        "",
        {"text": "立即扫码预约", "size": 24, "bold": True, "color": GREEN},
    ], page_num=84)

    # 第85页 - 尾页
    slide = add_title_slide(prs, "让您的下一个非标项目", "从救火走向掌控")

    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(2))
    tf = contact_box.text_frame
    tf.word_wrap = True

    contacts = [
        "服务热线：400-XXX-XXXX",
        "邮箱：sales@company.com",
        "官网：www.company.com",
        "地址：XX市XX区XX路XX号",
    ]

    for i, contact in enumerate(contacts):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = contact
        p.font.size = Pt(16)
        p.font.color.rgb = SILVER
        p.alignment = PP_ALIGN.CENTER

    # 保存文件
    output_path = "/Users/flw/non-standard-automation-pm/非标自动化项目管理系统_完整介绍.pptx"
    prs.save(output_path)
    print(f"PPT已生成: {output_path}")
    print(f"总页数: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    create_full_presentation()

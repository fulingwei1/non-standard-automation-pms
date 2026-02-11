# -*- coding: utf-8 -*-
"""Tests for app.services.ppt_generator.base_builder"""

import unittest
from unittest.mock import MagicMock, patch

from pptx import Presentation

from app.services.ppt_generator.base_builder import BaseSlideBuilder


class TestBaseSlideBuilder(unittest.TestCase):

    def setUp(self):
        self.prs = Presentation()
        self.builder = BaseSlideBuilder(self.prs)

    def test_init(self):
        self.assertIs(self.builder.prs, self.prs)
        self.assertIsNotNone(self.builder.config)

    def test_add_title_slide_no_subtitle(self):
        slide = self.builder.add_title_slide("Test Title")
        self.assertIsNotNone(slide)
        self.assertEqual(len(self.prs.slides), 1)

    def test_add_title_slide_with_subtitle(self):
        slide = self.builder.add_title_slide("Title", "Subtitle")
        self.assertIsNotNone(slide)
        # Should have background + title + subtitle = 3 shapes
        self.assertGreaterEqual(len(slide.shapes), 3)

    def test_add_section_slide_no_subtitle(self):
        slide = self.builder.add_section_slide("Section 1")
        self.assertIsNotNone(slide)

    def test_add_section_slide_with_subtitle(self):
        slide = self.builder.add_section_slide("Section 1", "Description")
        self.assertIsNotNone(slide)
        self.assertGreaterEqual(len(slide.shapes), 4)

    def test_add_white_background(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = self.builder._add_white_background(slide)
        self.assertIsNotNone(bg)

    def test_add_top_bar(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bar = self.builder._add_top_bar(slide)
        self.assertIsNotNone(bar)

    def test_add_slide_title(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        box = self.builder._add_slide_title(slide, "Page Title")
        self.assertIsNotNone(box)

    def test_add_page_number(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        box = self.builder._add_page_number(slide, 5)
        self.assertIsNotNone(box)
        self.assertEqual(box.text_frame.paragraphs[0].text, "5")


if __name__ == "__main__":
    unittest.main()

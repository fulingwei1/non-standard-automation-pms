# -*- coding: utf-8 -*-
"""
第三十五批 - content_builder.py 单元测试
"""
import pytest

pytest.importorskip("pptx", reason="python-pptx not available")

try:
    from unittest.mock import MagicMock, patch, PropertyMock
    from pptx.util import Inches, Pt
    from app.services.ppt_generator.content_builder import ContentSlideBuilder
    IMPORT_OK = True
except Exception:
    IMPORT_OK = False

pytestmark = pytest.mark.skipif(not IMPORT_OK, reason="导入失败，跳过")


def make_builder():
    """创建 ContentSlideBuilder mock 实例"""
    mock_config = MagicMock()
    mock_config.DARK_BLUE = MagicMock()
    mock_config.CONTENT_FONT_SIZE = Pt(16)
    mock_prs = MagicMock()

    # 构造 mock slide
    mock_slide = MagicMock()
    mock_shape = MagicMock()
    mock_tf = MagicMock()
    mock_para = MagicMock()
    mock_tf.paragraphs = [mock_para]
    mock_shape.text_frame = mock_tf
    mock_slide.shapes.add_textbox.return_value = mock_shape
    mock_prs.slides.add_slide.return_value = mock_slide
    mock_prs.slide_layouts = [MagicMock() for _ in range(10)]

    builder = ContentSlideBuilder.__new__(ContentSlideBuilder)
    builder.prs = mock_prs
    builder.config = mock_config
    # mock 基类方法
    builder._add_white_background = MagicMock()
    builder._add_top_bar = MagicMock()
    builder._add_slide_title = MagicMock()
    builder._add_page_number = MagicMock()
    builder._fill_textbox = MagicMock()
    return builder, mock_slide


@pytest.mark.skipif(not IMPORT_OK, reason="导入失败")
class TestContentSlideBuilder:

    def test_add_content_slide_string_items(self):
        """纯字符串列表能正常添加"""
        builder, mock_slide = make_builder()
        result = builder.add_content_slide("标题", ["行一", "行二"])
        assert result == mock_slide
        builder._add_white_background.assert_called_once()
        builder._add_slide_title.assert_called_once_with(mock_slide, "标题")

    def test_add_content_slide_dict_items(self):
        """字典格式内容项能正常处理"""
        builder, mock_slide = make_builder()
        items = [
            {"text": "第一行", "size": 20, "bold": True},
            {"text": "第二行", "size": 16, "bold": False, "color": MagicMock()},
        ]
        result = builder.add_content_slide("测试标题", items)
        assert result == mock_slide

    def test_add_content_slide_with_page_num(self):
        """有页码时调用 _add_page_number"""
        builder, mock_slide = make_builder()
        builder.add_content_slide("标题", ["内容"], page_num=3)
        builder._add_page_number.assert_called_once_with(mock_slide, 3)

    def test_add_content_slide_without_page_num(self):
        """无页码时不调用 _add_page_number"""
        builder, mock_slide = make_builder()
        builder.add_content_slide("标题", ["内容"])
        builder._add_page_number.assert_not_called()

    def test_add_two_column_slide_calls_fill(self):
        """两栏幻灯片调用 _fill_textbox 两次"""
        builder, mock_slide = make_builder()
        left = ["左一", "左二"]
        right = ["右一"]
        builder.add_two_column_slide("两栏标题", left, right)
        assert builder._fill_textbox.call_count == 2

    def test_add_two_column_slide_with_page_num(self):
        """两栏幻灯片带页码"""
        builder, mock_slide = make_builder()
        builder.add_two_column_slide("标题", ["左"], ["右"], page_num=5)
        builder._add_page_number.assert_called_once_with(mock_slide, 5)

    def test_fill_textbox_sets_text(self):
        """_fill_textbox 为真实实现时正常填充"""
        builder, _ = make_builder()
        # 恢复真实方法
        builder._fill_textbox = ContentSlideBuilder._fill_textbox.__get__(builder)

        mock_tf = MagicMock()
        mock_para = MagicMock()
        mock_tf.paragraphs = [mock_para]
        mock_textbox = MagicMock()
        mock_textbox.text_frame = mock_tf

        builder._fill_textbox(mock_textbox, ["内容A", "内容B"])
        assert mock_para.text == "内容A"

    def test_add_content_slide_empty_list(self):
        """空内容列表不抛出异常"""
        builder, mock_slide = make_builder()
        result = builder.add_content_slide("空标题", [])
        assert result == mock_slide

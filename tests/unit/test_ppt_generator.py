# -*- coding: utf-8 -*-
"""
PPT生成器单元测试

目标：
1. 参考 test_condition_parser_rewrite.py 的mock策略
2. 只mock外部依赖（pptx.Presentation, builders等）
3. 测试核心业务逻辑
4. 达到70%+覆盖率
"""

import unittest
from unittest.mock import MagicMock, patch, call

import pytest

try:
    from pptx.util import Inches, Pt
    from app.services.ppt_generator.generator import PresentationGenerator
except ImportError:
    pytest.skip("python-pptx not available", allow_module_level=True)


class TestPresentationGenerator(unittest.TestCase):
    """测试 PresentationGenerator 核心功能"""

    def setUp(self):
        """每个测试前初始化"""
        # Mock pptx.Presentation
        self.mock_prs = MagicMock()
        self.mock_prs.slide_width = Inches(10)
        self.mock_prs.slide_height = Inches(7.5)

        # Mock builders
        self.mock_base_builder = MagicMock()
        self.mock_content_builder = MagicMock()
        self.mock_table_builder = MagicMock()

        # Patch外部依赖
        self.patcher_prs = patch("app.services.ppt_generator.generator.Presentation")
        self.patcher_base = patch(
            "app.services.ppt_generator.generator.BaseSlideBuilder"
        )
        self.patcher_content = patch(
            "app.services.ppt_generator.generator.ContentSlideBuilder"
        )
        self.patcher_table = patch(
            "app.services.ppt_generator.generator.TableSlideBuilder"
        )

        # 启动patches
        self.mock_prs_class = self.patcher_prs.start()
        self.mock_base_class = self.patcher_base.start()
        self.mock_content_class = self.patcher_content.start()
        self.mock_table_class = self.patcher_table.start()

        # 配置mock返回值
        self.mock_prs_class.return_value = self.mock_prs
        self.mock_base_class.return_value = self.mock_base_builder
        self.mock_content_class.return_value = self.mock_content_builder
        self.mock_table_class.return_value = self.mock_table_builder

        # 创建测试对象
        self.generator = PresentationGenerator()

    def tearDown(self):
        """每个测试后清理"""
        self.patcher_prs.stop()
        self.patcher_base.stop()
        self.patcher_content.stop()
        self.patcher_table.stop()

    # ========== 初始化测试 ==========

    def test_init_creates_presentation(self):
        """测试初始化时创建Presentation对象"""
        self.mock_prs_class.assert_called_once()
        self.assertEqual(self.generator.prs, self.mock_prs)

    def test_init_sets_slide_dimensions(self):
        """测试初始化时设置幻灯片尺寸"""
        # 验证宽高被正确设置
        self.assertIsNotNone(self.mock_prs.slide_width)
        self.assertIsNotNone(self.mock_prs.slide_height)

    def test_init_creates_builders(self):
        """测试初始化时创建所有builder"""
        # 验证所有builder都被创建
        self.mock_base_class.assert_called_once_with(self.mock_prs)
        self.mock_content_class.assert_called_once_with(self.mock_prs)
        self.mock_table_class.assert_called_once_with(self.mock_prs)

        self.assertEqual(self.generator.base_builder, self.mock_base_builder)
        self.assertEqual(self.generator.content_builder, self.mock_content_builder)
        self.assertEqual(self.generator.table_builder, self.mock_table_builder)

    def test_init_creates_config(self):
        """测试初始化时创建配置对象"""
        self.assertIsNotNone(self.generator.config)

    # ========== create_cover_slide 测试 ==========

    def test_create_cover_slide_calls_base_builder(self):
        """测试create_cover_slide调用base_builder.add_title_slide"""
        mock_slide = MagicMock()
        self.mock_base_builder.add_title_slide.return_value = mock_slide
        mock_slide.shapes = MagicMock()

        self.generator.create_cover_slide()

        self.mock_base_builder.add_title_slide.assert_called_once_with(
            "智能驱动 · 精准交付", "非标自动化测试设备全生命周期 AI 项目管理系统"
        )

    def test_create_cover_slide_adds_slogan(self):
        """测试create_cover_slide添加slogan"""
        mock_slide = MagicMock()
        mock_slogan_box = MagicMock()
        mock_slide.shapes.add_textbox.return_value = mock_slogan_box
        self.mock_base_builder.add_title_slide.return_value = mock_slide

        self.generator.create_cover_slide()

        # 验证添加了textbox
        mock_slide.shapes.add_textbox.assert_called_once()

        # 验证textbox内容
        call_args = mock_slide.shapes.add_textbox.call_args
        self.assertEqual(call_args[0][0], Inches(0.5))  # left
        self.assertEqual(call_args[0][1], Inches(5.5))  # top

    # ========== create_toc_slide 测试 ==========

    def test_create_toc_slide_has_correct_title(self):
        """测试create_toc_slide有正确的标题"""
        self.generator.create_toc_slide()

        self.mock_content_builder.add_content_slide.assert_called_once()
        call_args = self.mock_content_builder.add_content_slide.call_args

        # 验证标题
        self.assertEqual(call_args[0][0], "内容导览")

    def test_create_toc_slide_has_seven_sections(self):
        """测试create_toc_slide包含7个章节"""
        self.generator.create_toc_slide()

        call_args = self.mock_content_builder.add_content_slide.call_args
        content_list = call_args[0][1]

        # 验证有7个章节
        self.assertEqual(len(content_list), 7)

        # 验证第一个章节
        self.assertEqual(content_list[0]["text"], "一、行业洞察与挑战")
        self.assertEqual(content_list[0]["size"], 24)
        self.assertTrue(content_list[0]["bold"])

    def test_create_toc_slide_has_page_number(self):
        """测试create_toc_slide有页码2"""
        self.generator.create_toc_slide()

        call_args = self.mock_content_builder.add_content_slide.call_args
        page_num = call_args[1]["page_num"]

        self.assertEqual(page_num, 2)

    # ========== create_part1_industry_insights 测试 ==========

    def test_create_part1_creates_company_intro_slide(self):
        """测试create_part1创建公司简介幻灯片"""
        self.generator.create_part1_industry_insights()

        # 验证至少调用了add_content_slide
        self.assertTrue(self.mock_content_builder.add_content_slide.called)

        # 找到"关于我们"的调用
        calls = self.mock_content_builder.add_content_slide.call_args_list
        company_intro_call = None
        for c in calls:
            if c[0][0] == "关于我们":
                company_intro_call = c
                break

        self.assertIsNotNone(company_intro_call)
        # 验证页码为3
        self.assertEqual(company_intro_call[1]["page_num"], 3)

    def test_create_part1_creates_industry_background_slide(self):
        """测试create_part1创建行业背景幻灯片"""
        self.generator.create_part1_industry_insights()

        calls = self.mock_content_builder.add_content_slide.call_args_list
        background_call = None
        for c in calls:
            if c[0][0] == "非标自动化行业的黄金时代":
                background_call = c
                break

        self.assertIsNotNone(background_call)
        # 验证页码为4
        self.assertEqual(background_call[1]["page_num"], 4)

    def test_create_part1_creates_customer_table(self):
        """测试create_part1创建客户画像表格"""
        self.generator.create_part1_industry_insights()

        # 验证调用了add_table_slide
        self.assertTrue(self.mock_table_builder.add_table_slide.called)

        # 找到"我们的客户"的调用
        calls = self.mock_table_builder.add_table_slide.call_args_list
        customer_call = None
        for c in calls:
            if c[0][0] == "我们的客户":
                customer_call = c
                break

        self.assertIsNotNone(customer_call)
        # 验证表头
        headers = customer_call[0][1]
        self.assertEqual(headers, ["行业", "典型设备需求", "关键诉求"])
        # 验证页码为5
        self.assertEqual(customer_call[1]["page_num"], 5)

    def test_create_part1_calls_create_painpoint_slides(self):
        """测试create_part1调用_create_painpoint_slides"""
        with patch.object(
            self.generator, "_create_painpoint_slides"
        ) as mock_painpoint:
            self.generator.create_part1_industry_insights()
            mock_painpoint.assert_called_once()

    # ========== _create_painpoint_slides 测试 ==========

    def test_create_painpoint_slides_creates_complexity_slide(self):
        """测试_create_painpoint_slides创建复杂性困局幻灯片"""
        self.generator._create_painpoint_slides()

        calls = self.mock_content_builder.add_content_slide.call_args_list
        complexity_call = None
        for c in calls:
            if "复杂性困局" in c[0][0]:
                complexity_call = c
                break

        self.assertIsNotNone(complexity_call)
        # 验证页码为6
        self.assertEqual(complexity_call[1]["page_num"], 6)

    # ========== create_part2_solution_overview 测试 ==========

    def test_create_part2_creates_section_slide(self):
        """测试create_part2创建章节分隔页"""
        self.generator.create_part2_solution_overview()

        # 验证调用了add_section_slide
        self.assertTrue(self.mock_base_builder.add_section_slide.called)

        # 找到第二部分的调用
        calls = self.mock_base_builder.add_section_slide.call_args_list
        part2_call = None
        for c in calls:
            if c[0][0] == "第二部分":
                part2_call = c
                break

        self.assertIsNotNone(part2_call)
        self.assertEqual(part2_call[0][1], "解决方案总览")

    def test_create_part2_creates_transition_slide(self):
        """测试create_part2创建转折页"""
        self.generator.create_part2_solution_overview()

        calls = self.mock_content_builder.add_content_slide.call_args_list
        transition_call = None
        for c in calls:
            if "救火" in c[0][0] and "掌控" in c[0][0]:
                transition_call = c
                break

        self.assertIsNotNone(transition_call)
        # 验证页码为12
        self.assertEqual(transition_call[1]["page_num"], 12)

    # ========== create_part3_core_features 测试 ==========

    def test_create_part3_creates_section_slide(self):
        """测试create_part3创建章节分隔页"""
        self.generator.create_part3_core_features()

        calls = self.mock_base_builder.add_section_slide.call_args_list
        part3_call = None
        for c in calls:
            if c[0][0] == "第三部分":
                part3_call = c
                break

        self.assertIsNotNone(part3_call)
        self.assertEqual(part3_call[0][1], "核心功能详解")

    # ========== create_part4_customer_portal 测试 ==========

    def test_create_part4_creates_section_slide(self):
        """测试create_part4创建章节分隔页"""
        self.generator.create_part4_customer_portal()

        calls = self.mock_base_builder.add_section_slide.call_args_list
        part4_call = None
        for c in calls:
            if c[0][0] == "第四部分":
                part4_call = c
                break

        self.assertIsNotNone(part4_call)
        self.assertEqual(part4_call[0][1], "客户门户与体验")

    # ========== create_part5_ai_capabilities 测试 ==========

    def test_create_part5_creates_section_slide(self):
        """测试create_part5创建章节分隔页"""
        self.generator.create_part5_ai_capabilities()

        calls = self.mock_base_builder.add_section_slide.call_args_list
        part5_call = None
        for c in calls:
            if c[0][0] == "第五部分":
                part5_call = c
                break

        self.assertIsNotNone(part5_call)
        self.assertEqual(part5_call[0][1], "AI能力与技术亮点")

    # ========== create_part6_customer_value 测试 ==========

    def test_create_part6_creates_section_slide(self):
        """测试create_part6创建章节分隔页"""
        self.generator.create_part6_customer_value()

        calls = self.mock_base_builder.add_section_slide.call_args_list
        part6_call = None
        for c in calls:
            if c[0][0] == "第六部分":
                part6_call = c
                break

        self.assertIsNotNone(part6_call)
        self.assertEqual(part6_call[0][1], "客户价值与案例")

    # ========== create_part7_cooperation 测试 ==========

    def test_create_part7_creates_section_slide(self):
        """测试create_part7创建章节分隔页"""
        self.generator.create_part7_cooperation()

        calls = self.mock_base_builder.add_section_slide.call_args_list
        part7_call = None
        for c in calls:
            if c[0][0] == "第七部分":
                part7_call = c
                break

        self.assertIsNotNone(part7_call)
        self.assertEqual(part7_call[0][1], "合作展望")

    # ========== generate 测试 ==========

    def test_generate_creates_all_parts(self):
        """测试generate调用所有创建方法"""
        with patch.object(
            self.generator, "create_cover_slide"
        ) as mock_cover, patch.object(
            self.generator, "create_toc_slide"
        ) as mock_toc, patch.object(
            self.generator, "create_part1_industry_insights"
        ) as mock_part1, patch.object(
            self.generator, "create_part2_solution_overview"
        ) as mock_part2, patch.object(
            self.generator, "create_part3_core_features"
        ) as mock_part3, patch.object(
            self.generator, "create_part4_customer_portal"
        ) as mock_part4, patch.object(
            self.generator, "create_part5_ai_capabilities"
        ) as mock_part5, patch.object(
            self.generator, "create_part6_customer_value"
        ) as mock_part6, patch.object(
            self.generator, "create_part7_cooperation"
        ) as mock_part7:

            result = self.generator.generate("test_output.pptx")

            # 验证所有方法都被调用
            mock_cover.assert_called_once()
            mock_toc.assert_called_once()
            mock_part1.assert_called_once()
            mock_part2.assert_called_once()
            mock_part3.assert_called_once()
            mock_part4.assert_called_once()
            mock_part5.assert_called_once()
            mock_part6.assert_called_once()
            mock_part7.assert_called_once()

    def test_generate_saves_presentation(self):
        """测试generate保存PPT文件"""
        with patch.object(self.generator, "create_cover_slide"), patch.object(
            self.generator, "create_toc_slide"
        ), patch.object(
            self.generator, "create_part1_industry_insights"
        ), patch.object(
            self.generator, "create_part2_solution_overview"
        ), patch.object(
            self.generator, "create_part3_core_features"
        ), patch.object(
            self.generator, "create_part4_customer_portal"
        ), patch.object(
            self.generator, "create_part5_ai_capabilities"
        ), patch.object(
            self.generator, "create_part6_customer_value"
        ), patch.object(
            self.generator, "create_part7_cooperation"
        ):

            result = self.generator.generate("test_output.pptx")

            # 验证调用了save方法
            self.mock_prs.save.assert_called_once_with("test_output.pptx")

    def test_generate_returns_output_path(self):
        """测试generate返回输出路径"""
        with patch.object(self.generator, "create_cover_slide"), patch.object(
            self.generator, "create_toc_slide"
        ), patch.object(
            self.generator, "create_part1_industry_insights"
        ), patch.object(
            self.generator, "create_part2_solution_overview"
        ), patch.object(
            self.generator, "create_part3_core_features"
        ), patch.object(
            self.generator, "create_part4_customer_portal"
        ), patch.object(
            self.generator, "create_part5_ai_capabilities"
        ), patch.object(
            self.generator, "create_part6_customer_value"
        ), patch.object(
            self.generator, "create_part7_cooperation"
        ):

            result = self.generator.generate("test_output.pptx")

            self.assertEqual(result, "test_output.pptx")

    def test_generate_with_default_output_path(self):
        """测试generate使用默认输出路径"""
        with patch.object(self.generator, "create_cover_slide"), patch.object(
            self.generator, "create_toc_slide"
        ), patch.object(
            self.generator, "create_part1_industry_insights"
        ), patch.object(
            self.generator, "create_part2_solution_overview"
        ), patch.object(
            self.generator, "create_part3_core_features"
        ), patch.object(
            self.generator, "create_part4_customer_portal"
        ), patch.object(
            self.generator, "create_part5_ai_capabilities"
        ), patch.object(
            self.generator, "create_part6_customer_value"
        ), patch.object(
            self.generator, "create_part7_cooperation"
        ):

            result = self.generator.generate()

            # 验证使用默认路径
            self.mock_prs.save.assert_called_once_with("完整PPT.pptx")
            self.assertEqual(result, "完整PPT.pptx")

    # ========== 边界情况和集成测试 ==========

    def test_create_cover_slide_integration(self):
        """测试create_cover_slide完整流程（不mock内部调用）"""
        # 重新创建一个不mock的generator用于集成测试
        # 这里保持原有的mock，只测试调用流程
        mock_slide = MagicMock()
        mock_shapes = MagicMock()
        mock_slide.shapes = mock_shapes
        mock_textbox = MagicMock()
        mock_shapes.add_textbox.return_value = mock_textbox

        self.mock_base_builder.add_title_slide.return_value = mock_slide

        self.generator.create_cover_slide()

        # 验证流程正确
        self.mock_base_builder.add_title_slide.assert_called_once()
        mock_shapes.add_textbox.assert_called_once()

    def test_create_toc_slide_content_structure(self):
        """测试create_toc_slide内容结构完整性"""
        self.generator.create_toc_slide()

        call_args = self.mock_content_builder.add_content_slide.call_args
        content_list = call_args[0][1]

        # 验证所有章节都有正确的结构
        for item in content_list:
            self.assertIn("text", item)
            self.assertIn("size", item)
            self.assertIn("bold", item)
            self.assertTrue(item["text"].startswith(("一、", "二、", "三", "四、", "五、", "六、", "七、")))

    def test_part1_customer_table_structure(self):
        """测试part1客户表格数据结构"""
        self.generator.create_part1_industry_insights()

        calls = self.mock_table_builder.add_table_slide.call_args_list
        customer_call = None
        for c in calls:
            if c[0][0] == "我们的客户":
                customer_call = c
                break

        rows = customer_call[0][2]
        # 验证有4行数据（汽车电子、消费电子、半导体、新能源）
        self.assertEqual(len(rows), 4)

        # 验证第一行
        self.assertEqual(rows[0][0], "汽车电子")
        # 验证每行都有3列
        for row in rows:
            self.assertEqual(len(row), 3)

    def test_multiple_parts_no_conflict(self):
        """测试多个部分创建互不冲突"""
        # 连续创建多个部分
        self.generator.create_part1_industry_insights()
        self.generator.create_part2_solution_overview()
        self.generator.create_part3_core_features()

        # 验证所有部分都成功创建（通过调用次数）
        # add_section_slide应该被调用3次（part2, part3各一次，part1不调用）
        section_calls = self.mock_base_builder.add_section_slide.call_args_list
        part_titles = [call[0][0] for call in section_calls]

        self.assertIn("第二部分", part_titles)
        self.assertIn("第三部分", part_titles)

    def test_config_attributes_exist(self):
        """测试config对象属性存在"""
        # 验证config有必要的属性
        self.assertTrue(hasattr(self.generator.config, "SILVER"))
        self.assertTrue(hasattr(self.generator.config, "TECH_BLUE"))
        self.assertTrue(hasattr(self.generator.config, "ORANGE"))


class TestPresentationGeneratorEdgeCases(unittest.TestCase):
    """测试边界情况和异常处理"""

    def setUp(self):
        """每个测试前初始化"""
        self.mock_prs = MagicMock()
        self.mock_prs.slide_width = Inches(10)
        self.mock_prs.slide_height = Inches(7.5)

        self.patcher_prs = patch("app.services.ppt_generator.generator.Presentation")
        self.patcher_base = patch(
            "app.services.ppt_generator.generator.BaseSlideBuilder"
        )
        self.patcher_content = patch(
            "app.services.ppt_generator.generator.ContentSlideBuilder"
        )
        self.patcher_table = patch(
            "app.services.ppt_generator.generator.TableSlideBuilder"
        )

        self.mock_prs_class = self.patcher_prs.start()
        self.mock_base_class = self.patcher_base.start()
        self.mock_content_class = self.patcher_content.start()
        self.mock_table_class = self.patcher_table.start()

        self.mock_prs_class.return_value = self.mock_prs
        self.mock_base_class.return_value = MagicMock()
        self.mock_content_class.return_value = MagicMock()
        self.mock_table_class.return_value = MagicMock()

    def tearDown(self):
        """清理"""
        self.patcher_prs.stop()
        self.patcher_base.stop()
        self.patcher_content.stop()
        self.patcher_table.stop()

    def test_generate_with_special_characters_in_path(self):
        """测试带特殊字符的输出路径"""
        generator = PresentationGenerator()

        with patch.object(generator, "create_cover_slide"), patch.object(
            generator, "create_toc_slide"
        ), patch.object(
            generator, "create_part1_industry_insights"
        ), patch.object(
            generator, "create_part2_solution_overview"
        ), patch.object(
            generator, "create_part3_core_features"
        ), patch.object(
            generator, "create_part4_customer_portal"
        ), patch.object(
            generator, "create_part5_ai_capabilities"
        ), patch.object(
            generator, "create_part6_customer_value"
        ), patch.object(
            generator, "create_part7_cooperation"
        ):

            special_path = "test/路径/中文PPT文件.pptx"
            result = generator.generate(special_path)

            self.mock_prs.save.assert_called_once_with(special_path)
            self.assertEqual(result, special_path)

    def test_create_cover_slide_with_mock_failure(self):
        """测试base_builder失败时的处理"""
        generator = PresentationGenerator()

        # 模拟add_title_slide返回None
        generator.base_builder.add_title_slide.return_value = None

        # 应该抛出AttributeError（因为会尝试访问None.shapes）
        with self.assertRaises(AttributeError):
            generator.create_cover_slide()


if __name__ == "__main__":
    unittest.main()

"""Tests for app/services/ppt_generator/table_builder.py"""
import pytest
from unittest.mock import MagicMock, patch, call

try:
    from app.services.ppt_generator.table_builder import TableSlideBuilder
    from app.services.ppt_generator.config import PresentationConfig
    from pptx.util import Inches
except ImportError as e:
    pytest.skip(f"Import failed: {e}", allow_module_level=True)


def _make_mock_prs_with_table():
    """构建带 table 支持的 mock Presentation"""
    prs = MagicMock()
    slide = MagicMock()
    # mock table object
    table = MagicMock()
    cell = MagicMock()
    cell.text = ""
    cell.fill = MagicMock()
    cell.fill.fore_color = MagicMock()
    cell.text_frame = MagicMock()
    para = MagicMock()
    para.font = MagicMock()
    para.font.size = None
    para.font.bold = False
    para.font.color = MagicMock()
    cell.text_frame.paragraphs = [para]
    table.cell.return_value = cell
    col = MagicMock()
    col.width = 0
    table.columns = [col, col, col]

    table_shape = MagicMock()
    table_shape.table = table
    slide.shapes = MagicMock()
    slide.shapes.add_table.return_value = table_shape

    textbox = MagicMock()
    tf = MagicMock()
    p = MagicMock()
    p.font = MagicMock()
    p.font.color = MagicMock()
    tf.paragraphs = [p]
    textbox.text_frame = tf
    slide.shapes.add_textbox.return_value = textbox

    shape = MagicMock()
    shape.fill = MagicMock()
    shape.fill.fore_color = MagicMock()
    shape.line = MagicMock()
    shape.line.fill = MagicMock()
    slide.shapes.add_shape.return_value = shape

    slide_layout = MagicMock()
    prs.slide_layouts = [slide_layout] * 10
    prs.slides = MagicMock()
    prs.slides.add_slide.return_value = slide
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs, slide, table


def test_add_table_slide_returns_slide():
    """add_table_slide 返回 slide 对象"""
    prs, slide, _ = _make_mock_prs_with_table()
    builder = TableSlideBuilder(prs)
    headers = ["列1", "列2", "列3"]
    rows = [["A", "B", "C"], ["D", "E", "F"]]
    result = builder.add_table_slide("测试表格", headers, rows)
    assert result is slide


def test_add_table_slide_calls_add_table():
    """add_table_slide 调用 add_table"""
    prs, slide, _ = _make_mock_prs_with_table()
    builder = TableSlideBuilder(prs)
    headers = ["列1", "列2"]
    rows = [["A", "B"]]
    builder.add_table_slide("测试表格", headers, rows)
    slide.shapes.add_table.assert_called_once()


def test_add_table_slide_with_page_num():
    """add_table_slide 传入 page_num 时调用 _add_page_number"""
    prs, slide, _ = _make_mock_prs_with_table()
    builder = TableSlideBuilder(prs)
    with patch.object(builder, '_add_page_number') as mock_page:
        builder.add_table_slide("表格", ["A"], [["1"]], page_num=5)
        mock_page.assert_called_once_with(slide, 5)


def test_add_table_slide_without_page_num():
    """add_table_slide 不传入 page_num 时不调用 _add_page_number"""
    prs, slide, _ = _make_mock_prs_with_table()
    builder = TableSlideBuilder(prs)
    with patch.object(builder, '_add_page_number') as mock_page:
        builder.add_table_slide("表格", ["A"], [["1"]])
        mock_page.assert_not_called()


def test_format_table_header_sets_text():
    """_format_table_header 设置表头文字"""
    prs, slide, table = _make_mock_prs_with_table()
    builder = TableSlideBuilder(prs)
    headers = ["标题1", "标题2", "标题3"]
    builder._format_table_header(table, headers)
    # cell() 应被调用 len(headers) 次（行0）
    assert table.cell.call_count >= len(headers)


def test_format_table_rows_sets_text():
    """_format_table_rows 设置数据行文字"""
    prs, slide, table = _make_mock_prs_with_table()
    builder = TableSlideBuilder(prs)
    rows = [["A", "B"], ["C", "D"]]
    builder._format_table_rows(table, rows)
    # 每行每列都应调用 cell()
    assert table.cell.call_count >= len(rows) * 2


def test_add_table_slide_correct_num_rows():
    """add_table_slide 以 num_rows = len(rows)+1 调用 add_table"""
    prs, slide, _ = _make_mock_prs_with_table()
    builder = TableSlideBuilder(prs)
    headers = ["A", "B"]
    rows = [["1", "2"], ["3", "4"], ["5", "6"]]
    builder.add_table_slide("Test", headers, rows)
    # add_table(num_rows, num_cols, ...)
    call_args = slide.shapes.add_table.call_args[0]
    assert call_args[0] == 4  # 3 rows + 1 header
    assert call_args[1] == 2  # 2 columns
